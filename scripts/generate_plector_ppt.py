#!/usr/bin/env python3
"""Plector项目PPT生成脚本"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
TITLE_COLOR = RgbColor(0, 102, 204)
ACCENT_COLOR = RgbColor(0, 153, 76)
DARK_BG = RgbColor(30, 30, 46)
LIGHT_TEXT = RgbColor(255, 255, 255)
GRAY_TEXT = RgbColor(100, 100, 100)


def add_title_slide(prs, title, subtitle):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(13.333), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = LIGHT_TEXT
    p2.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, bullets):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BG
    header.line.fill.background()
    
    # 标题文字
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    
    # 内容区域
    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.933), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_before = Pt(12)
        p.space_after = Pt(6)
    
    return slide


def add_two_column_slide(prs, title, left_content, right_content):
    """添加双栏幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BG
    header.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    
    # 左栏
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf_left.paragraphs[0]
        else:
            p = tf_left.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_before = Pt(8)
    
    # 右栏
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf_right.paragraphs[0]
        else:
            p = tf_right.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_before = Pt(8)
    
    return slide


def add_architecture_slide(prs, title, modules):
    """添加架构图幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BG
    header.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    
    # 绘制模块框
    colors = [DARK_BG, RgbColor(0, 128, 128), RgbColor(128, 0, 128), 
              RgbColor(200, 100, 50), RgbColor(50, 100, 150)]
    
    positions = [
        (Inches(0.5), Inches(1.5), Inches(3.5), Inches(1.2)),   # core
        (Inches(4.5), Inches(1.5), Inches(2.5), Inches(1.2)),   # skills
        (Inches(7.5), Inches(1.5), Inches(2.5), Inches(1.2)),   # tools
        (Inches(10.5), Inches(1.5), Inches(2.5), Inches(1.2)),  # channels
        (Inches(7.5), Inches(3.0), Inches(2.5), Inches(1.2)),   # servers
    ]
    
    for i, (name, desc) in enumerate(modules[:5]):
        left, top, width, height = positions[i]
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()
        
        txBox = slide.shapes.add_textbox(left, top + Inches(0.3), width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = LIGHT_TEXT
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = LIGHT_TEXT
        p2.alignment = PP_ALIGN.CENTER
    
    # 添加模块列表
    details = [
        ("core/", "agent_loop, closure_engine, skill_registry, event_bus, context_builder..."),
        ("skills/", "agency_orchestrator, auto_developer, code_writer, memory, web_search..."),
        ("tools/", "49个工具函数，无状态无治理"),
        ("channels/", "cli.py, websocket.py, http.py"),
        ("servers/", "MCP Servers: agency-orchestrator, filesystem, sqlite..."),
    ]
    
    detail_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(2.8))
    tf_det = detail_box.text_frame
    tf_det.word_wrap = True
    for i, (name, desc) in enumerate(details):
        if i == 0:
            p = tf_det.paragraphs[0]
        else:
            p = tf_det.add_paragraph()
        p.text = f"{name}: {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(80, 80, 80)
        p.space_before = Pt(4)
    
    return slide


def add_feature_slide(prs, title, features):
    """添加特性卡片幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BG
    header.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    
    # 特性卡片
    card_colors = [DARK_BG, RgbColor(0, 128, 128), RgbColor(128, 0, 128), 
                   RgbColor(200, 100, 50), ACCENT_COLOR, RgbColor(100, 50, 150),
                   RgbColor(50, 100, 100)]
    
    card_width = Inches(3.7)
    card_height = Inches(2.5)
    start_x = Inches(0.5)
    start_y = Inches(1.6)
    gap = Inches(0.25)
    
    for i, (feature, desc) in enumerate(features[:6]):
        row = i // 3
        col = i % 3
        left = start_x + col * (card_width + gap)
        top = start_y + row * (card_height + gap)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = card_colors[i % len(card_colors)]
        card.line.fill.background()
        
        txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), card_width - Inches(0.3), card_height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = feature
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = LIGHT_TEXT
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = LIGHT_TEXT
    
    return slide


# ==================== 生成PPT内容 ====================

# 1. 封面
add_title_slide(prs, "Plector", "事件驱动的 AI Agent 引擎 | v1.8.0")

# 2. 项目概述
add_content_slide(prs, "项目概述", [
    "当前版本: v1.8.0 (已定稿)",
    "核心引擎: 13 个核心模块 | 技能系统: 9 个技能 | 工具集: 49 个工具",
    "架构风格: 事件驱动架构，组件异步解耦",
    "核心特性: ReAct自主决策 | 多LLM后端 | MCP协议 | 闭环引擎",
    "设计原则: 模块解耦 | 技能治理 | 闭环自愈",
    "代码质量: 单元测试覆盖率 ≥ 80% | Harness 7项自动化检查",
])

# 3. 核心能力
add_feature_slide(prs, "核心能力", [
    ("自主决策 (ReAct)", "LLM推理 → 调用工具 → 观察 → 迭代的循环机制"),
    ("多LLM后端", "支持 Ollama / OpenAI / Anthropic 三大后端"),
    ("技能系统", "插件化技能，MCP格式定义，最多15个技能"),
    ("事件驱动", "CloudEvents 1.0标准，组件异步入耦"),
    ("MCP协议", "连接外部MCP Server，引入现成工具"),
    ("闭环引擎", "条件图执行，支持自动修复"),
])

# 4. 项目架构
add_architecture_slide(prs, "项目架构", [
    ("core/", "核心引擎"),
    ("skills/", "技能系统"),
    ("tools/", "工具集"),
    ("channels/", "接入渠道"),
    ("servers/", "MCP Servers"),
])

# 5. 核心模块详解
add_two_column_slide(prs, "核心模块详解", [
    "Agent Loop (agent_loop.py)",
    "  - 实现ReAct模式",
    "  - 管理LLM调用、工具执行、结果回填",
    "  - 将技能也注册为工具统一调用",
    "",
    "Closure Engine (closure_engine.py)",
    "  - 条件图解析与执行",
    "  - 支持自动修复",
    "  - 配置统一路径: config/closed_loops.yaml",
    "",
    "Skill Registry (skill_registry.py)",
    "  - 技能注册与管理",
    "  - 技能元数据存储",
], [
    "Event Bus (event_bus.py)",
    "  - 事件总线",
    "  - 支持同步/异步事件处理",
    "  - CloudEvents 1.0兼容",
    "",
    "Context Builder (context_builder.py)",
    "  - 上下文构建",
    "  - 从.md文件加载配置",
    "  - 支持多角色配置(AGENTS/SOUL/USER)",
    "",
    "Governance (governance.py)",
    "  - 技能治理",
    "  - 健康分析、淘汰机制",
])

# 6. 技能系统
add_two_column_slide(prs, "技能系统 (9个技能)", [
    "agency_orchestrator",
    "  - 多智能体YAML工作流引擎",
    "  - 174个AI角色，支持DAG并行执行",
    "  - 支持Resume断点续跑",
    "",
    "auto_developer",
    "  - 一键自动开发流水线",
    "  - 需求→代码全自动",
    "",
    "code_writer",
    "  - 代码编写技能",
    "  - 支持写入、读取、修改代码文件",
    "",
    "test_runner",
    "  - 测试运行技能",
    "  - 支持pytest并返回结果",
], [
    "memory",
    "  - 记忆管理技能",
    "  - 存储对话历史、用户偏好、知识记忆",
    "",
    "web_search",
    "  - 网页搜索技能",
    "  - 使用博查API(国内可用)",
    "",
    "file_utils",
    "  - 文件操作技能",
    "  - 列表、复制、移动、删除文件",
    "",
    "health_monitor",
    "  - 系统健康监控",
    "  - CPU、内存、磁盘使用率",
    "",
    "error_knowledge",
    "  - 错误记录与分类",
    "  - 存储到本地知识库",
])

# 7. 技术栈
add_content_slide(prs, "技术栈与依赖", [
    "Python 3.11+ | asyncio异步编程",
    "LLM后端: Ollama (本地) | OpenAI API | Anthropic API",
    "MCP协议: Model Context Protocol 实现",
    "数据存储: SQLite + 向量存储 (vector_memory)",
    "事件标准: CloudEvents 1.0",
    "测试框架: pytest + pre-commit hooks",
    "代码规范: 单函数≤50行 | 返回值{success, data, error}",
])

# 8. 快速开始
add_content_slide(prs, "快速开始", [
    "# 克隆项目",
    "git clone https://github.com/biandeshen/Plector.git",
    "cd Plector",
    "",
    "# 配置LLM后端 (三选一)",
    "ollama pull qwen3:4b && ollama serve  # Ollama(本地)",
    "export OPENAI_API_KEY=\"sk-xxx\"     # OpenAI",
    "export ANTHROPIC_API_KEY=\"sk-ant-xxx\" # Anthropic",
    "",
    "# 运行",
    "python channels/cli.py --query \"你好\"      # CLI模式",
    "python channels/websocket.py --port 8080  # Web模式",
])

# 9. 设计原则
add_content_slide(prs, "设计原则", [
    "硬性规则:",
    "  1. core/ 不依赖 skills/ 和 tools/",
    "  2. 技能数量 ≤ 15 个",
    "  3. 函数不超过 50 行",
    "  4. 返回值格式: {\"success\", \"data\", \"error\"}",
    "",
    "Plector性格:",
    "  • 务实、简洁、高效",
    "  • 先验证再优化",
    "  • 出错时返回结构化错误，不抛异常",
    "  • 拒绝过度工程化",
])

# 10. 总结
add_title_slide(prs, "谢谢观看", "Plector - 让AI Agent开发更简单")

# 保存文件
output_path = "Plector/presentation/Plector_Introduction.pptx"
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
